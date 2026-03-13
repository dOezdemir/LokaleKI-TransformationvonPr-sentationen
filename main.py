# ==========================================================
# Phase 1 
# ==========================================================
import os
import io
import json
import hashlib
import logging
import platform
import time
import sys
import argparse
import re
from typing import Optional, Dict, List, Any
import pptx
import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image


# ===========================================================================================
# Ordner vorbereiten: Erstellt die benötigten Ausgabeordner, falls sie noch nicht existieren
# ===========================================================================================
os.makedirs("export/assets", exist_ok=True)
os.makedirs("data/raw/logs", exist_ok=True)
os.makedirs("data/processed", exist_ok=True)


# ============================================================================================================================
# Hashing-Funktion: Berechnet SHA-256-Hash für eine Datei damit später überprüft werden kann, ob isch eine Datei verändert hat
# ============================================================================================================================
def hash_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()

# Durchsucht ein Verzeichnis rekursiv nach Dateien und schreibt deren Hash-Werte ins Audit-Log
def audit_hash_artifacts(root_dir: str, prefix: str) -> None:
    files = []
    for base, _, names in os.walk(root_dir):
        for n in names:
            p = os.path.join(base, n)
            if os.path.isfile(p):
                files.append(p)

    for p in sorted(files):
        rel = os.path.realpath(p).replace("\\", "/")
        audit_logger.info(f"{prefix} | {rel} | Hash: {hash_file(p)}")


# =================================================================================
# Logging einrichten: Konfiguration des normalen Loggings für Prozessinformationen
# =================================================================================
logging.basicConfig(
    filename="data/raw/logs/pipeline.log",
    level=logging.INFO,
    format="%(levelname)s | %(asctime)s | %(message)s | %(funcName)s",
)

# Logger für Audit-Zwecke
audit_logger = logging.getLogger("audit")
audit_logger.propagate = False

# Audit-Log-Datei
audit_handler = logging.FileHandler("data/raw/logs/audit.log")
audit_handler.setFormatter(logging.Formatter("%(asctime)s | %(message)s"))

# Verhindert, dass derselbe Handler mehrfach hinzugefügt wird
if not any(
    isinstance(h, logging.FileHandler)
    and getattr(h, "baseFilename", "") == audit_handler.baseFilename
    for h in audit_logger.handlers
):
    audit_logger.addHandler(audit_handler)

audit_logger.setLevel(logging.INFO)


# ==============================================================================================================
# Helpers: Versucht einen Wert in float umzuwandeln und gibt None zurück, wenn die Umwandlung nicht möglich ist
# ==============================================================================================================
def safe_float(x: Any) -> Optional[float]:
    try:
        return float(x)
    except Exception:
        return None


# ==================================================================================================
# Helpers (PDF Farben / Ausrichtung): Wandelt einen PDF-Farbwert Integer in ein RGB-Dictionary um
# ==================================================================================================
def pdf_color_to_rgb(color_int: Optional[int]) -> Optional[Dict[str, int]]:
    if color_int is None:
        return None
    try:
        r = (int(color_int) >> 16) & 255
        g = (int(color_int) >> 8) & 255
        b = int(color_int) & 255
        return {"r": r, "g": g, "b": b}
    except Exception:
        return None


# Wandelt ein RGB-Dictionary in einen Hex-Farbwert um, z. B. #FF0000
def pdf_rgb_to_hex(rgb: Optional[Dict[str, int]]) -> Optional[str]:
    if not rgb:
        return None
    try:
        return "#{:02X}{:02X}{:02X}".format(int(rgb["r"]), int(rgb["g"]), int(rgb["b"]))
    except Exception:
        return None


# Schätzt die Textausrichtung in einem PDF anhand der Position auf der Seite
def infer_pdf_align(x0, x1, page_width, tol=0.06) -> str:
    if x0 is None or x1 is None or not page_width:
        return "left"

    left_margin = x0 / page_width
    right_margin = (page_width - x1) / page_width
    center_offset = abs(((x0 + x1) / 2) - (page_width / 2)) / page_width

    if center_offset < tol:
        return "center"
    if right_margin < tol and left_margin > tol:
        return "right"
    return "left"


# ==========================================================
# Helpers (PPTX Text / Farbe / Bullets)
# ==========================================================
def pptx_align_to_css(alignment) -> Optional[str]:

    if alignment is None:
        return None
    s = str(alignment).lower()
    if "center" in s:
        return "center"
    if "right" in s:
        return "right"
    if "justify" in s:
        return "justify"
    if "left" in s:
        return "left"
    return None

# Liest die Schriftfarbe eines Text-Runs aus PowerPoint aus, falls keine Farbe vorhanden ist, wird None zurückgegeben
def safe_pptx_rgb(run) -> Optional[str]:
    try:
        if run is None or run.font is None:
            return None
        c = run.font.color
        if c is None:
            return None
        rgb = c.rgb
        if rgb is None:
            return None
        return str(rgb)
    except Exception:
        return None

# Prüft per XML, ob ein Absatz in PowerPoint als Aufzählung formatiert ist
def pptx_paragraph_is_bullet(p) -> bool:
    try:
        xml = p._p.xml
        if "<a:buNone" in xml:
            return False
        return ("<a:buChar" in xml) or ("<a:buAutoNum" in xml) or ("<a:buBlip" in xml)
    except Exception:
        return False


# ==========================================================
# Diagramm-Region Detection Helpers (PDF Zeichnungen)
# ==========================================================
PERCENT_RE = re.compile(r"^\s*\d{1,3}\s*%\s*$")

# Bildet die kleinste gemeinsame Bounding Box aus zwei Rechtecken
def rect_union(a: List[float], b: List[float]) -> List[float]:
    x0 = min(a[0], b[0])
    y0 = min(a[1], b[1])
    x1 = max(a[2], b[2])
    y1 = max(a[3], b[3])
    return [x0, y0, x1, y1]

# Berechnet die Fläche eines Rechtecks
def rect_area(r: List[float]) -> float:
    return max(0.0, (r[2] - r[0])) * max(0.0, (r[3] - r[1]))

# Vergrößert ein Rechteck in alle Richtungen um pad
def rect_expand(r: List[float], pad: float) -> List[float]:
    return [r[0] - pad, r[1] - pad, r[2] + pad, r[3] + pad]

# Prüft, ob sich zwei Rechtecke schneiden
def rect_intersects(a: List[float], b: List[float]) -> bool:
    return not (a[2] < b[0] or a[0] > b[2] or a[3] < b[1] or a[1] > b[3])

# Berechnet den Abstand zwischen zwei Rechtecken
def rect_distance(a: List[float], b: List[float]) -> float:
    dx = max(0.0, max(b[0] - a[2], a[0] - b[2]))
    dy = max(0.0, max(b[1] - a[3], a[1] - b[3]))
    return (dx * dx + dy * dy) ** 0.5

# Fasst nahe beieinander liegende Rechtecke zu größeren Bereichen zusammen
def cluster_rects(rects: List[List[float]], max_dist: float) -> List[List[float]]:
    clusters: List[List[float]] = []
    used = [False] * len(rects)

    for i in range(len(rects)):
        if used[i]:
            continue
        used[i] = True
        cluster = [rects[i]]
        changed = True

        while changed:
            changed = False
            cluster_bbox = cluster[0]
            for r in cluster[1:]:
                cluster_bbox = rect_union(cluster_bbox, r)

            for j in range(len(rects)):
                if used[j]:
                    continue
                d = rect_distance(cluster_bbox, rects[j])
                if d <= max_dist or rect_intersects(cluster_bbox, rects[j]):
                    used[j] = True
                    cluster.append(rects[j])
                    changed = True

        bbox = cluster[0]
        for r in cluster[1:]:
            bbox = rect_union(bbox, r)
        clusters.append(bbox)

    return clusters

# Extrahiert Rechtecke aus PDF-Zeichnungen, die groß genug sind, um eventuell Diagrammteile zu sein
def extract_shape_rects_from_drawings(
    drawings, page_width, page_height, min_area=200.0
):
    rects = []
    page_area = page_width * page_height

    for d in drawings or []:
        r = d.get("rect")
        if not r:
            continue

        r2 = [float(r.x0), float(r.y0), float(r.x1), float(r.y1)]
        area = rect_area(r2)

        if area < min_area:
            continue
        if area > 0.8 * page_area:
            continue

        rects.append(r2)

    return rects


# Erkennt mögliche Diagrammregionen im PDF anhand gezeichneter Formen und Prozent-Labels in der Nähe
def detect_chart_regions(
    page, text_elements: List[dict], drawings: Optional[List[dict]]
) -> List[dict]:
    pw = float(page.rect.width)
    ph = float(page.rect.height)
    pad = 0.01 * max(pw, ph)
    max_dist = 0.03 * max(pw, ph)

    shape_rects = extract_shape_rects_from_drawings(drawings, pw, ph, min_area=300.0)
    if not shape_rects:
        return []

    clusters = cluster_rects(shape_rects, max_dist=max_dist)
    clusters = [c for c in clusters if rect_area(c) >= 0.01 * pw * ph]

    chart_regions = []
    for bbox in clusters:
        bbox2 = rect_expand(bbox, pad)

        labels = []
        for el in text_elements:
            if el.get("typ") != "text":
                continue
            txt = (el.get("inhalt") or "").strip()
            if not PERCENT_RE.match(txt):
                continue

            pos = el.get("position") or {}
            x = pos.get("links")
            y = pos.get("oben")
            w = pos.get("breite")
            h = pos.get("höhe")
            if x is None or y is None or w is None or h is None:
                continue

            tr = [float(x), float(y), float(x + w), float(y + h)]
            if rect_intersects(bbox2, tr) or rect_distance(bbox2, tr) <= pad:
                labels.append(
                    {
                        "text": txt,
                        "position": {
                            "links": float(x),
                            "oben": float(y),
                            "breite": float(w),
                            "höhe": float(h),
                        },
                    }
                )
        # Region wird als Diagramm gespeichert, wenn Labels gefunden wurden oder der Bereich groß genug ist
        if labels or rect_area(bbox) >= 0.04 * pw * ph:
            chart_regions.append({"bbox": bbox, "labels": labels})

    return chart_regions


# ==========================================================
# PDF Text Helpers
# ==========================================================
# Prüft, ob ein Text-Span im PDF fett formatiert ist
def pdf_span_is_bold(span: dict) -> Optional[bool]:
    try:
        flags = int(span.get("flags", 0))
        font = (span.get("font") or "").lower()

        is_bold_by_name = ("bold" in font) or ("black" in font) or ("heavy" in font)
        is_bold_by_flags = bool(flags & 16)

        return bool(is_bold_by_name or is_bold_by_flags)
    except Exception:
        return None

# Prüft, ob ein Text-Span im PDF kursiv formatiert ist
def pdf_span_is_italic(span: dict) -> Optional[bool]:
    try:
        flags = int(span.get("flags", 0))
        font = (span.get("font") or "").lower()

        is_italic_by_name = ("italic" in font) or ("oblique" in font)
        is_italic_by_flags = bool(flags & 2)

        return bool(is_italic_by_name or is_italic_by_flags)
    except Exception:
        return None

# Regulärer Ausdruck für Bullet-Points oder nummerierte Listen
BULLET_RE = re.compile(r"^\s*([•\-\u2022\u25CF\u25AA\u25E6]|(\d+[\.\)]))\s+")

# Prüft, ob eine Textzeile wie ein Listenpunkt aussieht
def pdf_line_is_bullet(line_text: str) -> bool:
    if not line_text:
        return False
    return bool(BULLET_RE.match(line_text))


# Liest Textblöcke aus einer PDF-Seite aus und speichert Formatierungs- und Positionsinformationen
def extract_pdf_text_blocks(page: "fitz.Page") -> List[dict]:
    out: List[dict] = []

    pw = float(page.rect.width)
    text_dict = page.get_text("dict")

    for b in text_dict.get("blocks", []):
        if b.get("type") != 0:
            continue

        bbox = b.get("bbox")
        if not bbox or len(bbox) != 4:
            continue

        x0, y0, x1, y1 = map(float, bbox)
        block_width = max(0.0, x1 - x0)
        block_height = max(0.0, y1 - y0)

        lines = b.get("lines", []) or []
        paragraphs_out = []
        full_text_lines = []

        for line in lines:
            spans = line.get("spans", []) or []
            runs_out = []
            line_text_parts = []

            for sp in spans:
                t = sp.get("text") or ""
                if t == "":
                    continue

                line_text_parts.append(t)

                rgb = pdf_color_to_rgb(sp.get("color"))
                runs_out.append(
                    {
                        "text": t,
                        "schriftgröße": safe_float(sp.get("size")),
                        "schriftart": sp.get("font"),
                        "farbe": pdf_rgb_to_hex(rgb),
                        "bold": pdf_span_is_bold(sp),
                        "italic": pdf_span_is_italic(sp),
                        "underline": None,
                    }
                )

            line_text = "".join(line_text_parts).strip()
            if not line_text and not runs_out:
                continue

            full_text_lines.append(line_text)

            paragraphs_out.append(
                {
                    "level": 0,
                    "is_bullet": pdf_line_is_bullet(line_text),
                    "textausrichtung": infer_pdf_align(x0, x1, pw),
                    "runs": runs_out,
                    "plain_text": line_text,
                }
            )

        full_text = "\n".join([t for t in full_text_lines if t]).strip()
        if not full_text:
            continue

        out.append(
            {
                "typ": "text",
                "inhalt": full_text,
                "paragraphs": paragraphs_out,
                "z_index": None,
                "position": {
                    "links": x0,
                    "oben": y0,
                    "breite": block_width,
                    "höhe": block_height,
                },
            }
        )

    return out

# Versucht, den Seitentitel im PDF zu finden
def detect_pdf_title_from_textblocks(
    text_blocks: List[dict], page_height: float
) -> Optional[str]:
    best = None

    top_limit = 0.25 * float(page_height)

    for el in text_blocks:
        pos = el.get("position") or {}
        y = safe_float(pos.get("oben")) or 0.0
        if y > top_limit:
            continue

        paragraphs = el.get("paragraphs") or []
        max_size = None
        for p in paragraphs:
            for r in p.get("runs") or []:
                sz = safe_float(r.get("schriftgröße"))
                if sz is None:
                    continue
                max_size = sz if (max_size is None or sz > max_size) else max_size

        if max_size is None:
            continue

        text = (el.get("inhalt") or "").strip()
        if not text:
            continue

        score = (max_size * 1000.0) - y
        if best is None or score > best[0]:
            title_line = text.splitlines()[0].strip()
            best = (score, title_line)

    return best[1] if best else None


# ====================================================================================================================
# PPTX PARSING: Verarbeitet eine PowerPoint-Datei und exportiert Folieninformationen als JSON sowie Bilder als Assets
# ====================================================================================================================
def parse_pptx(pptx_path: str) -> None:
    logging.info(f"Starte PPTX-Parsing: {pptx_path}")
    start_time = time.time()

    input_hash = hash_file(pptx_path)
    audit_logger.info(f"PPTX Input: {pptx_path} | Hash: {input_hash}")

    pres = Presentation(pptx_path)
    bilder_ordner = "export/assets"
    bild_zaehler = 0

    for i, slide in enumerate(pres.slides, start=1):
        logging.info(f"Verarbeite PPTX-Folie {i}")

        folien_daten = {
            "foliennummer": i,
            "überschrift": "kein Titel gefunden",
            "canvas": {
                "width": int(pres.slide_width),
                "height": int(pres.slide_height),
                "unit": "emu",
            },
            "elemente": [],
        }

        z_index = 0

        # Falls ein offizieller Folienslot-Titel existiert, wird dieser übernommen
        titel_shape = getattr(slide.shapes, "title", None)
        if titel_shape is not None and getattr(titel_shape, "text", "").strip():
            folien_daten["überschrift"] = titel_shape.text.strip()

        for shape in slide.shapes:
            z_index += 1

            # Bildobjekte werden extrahiert und als Datei gespeichert
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                bild_zaehler += 1
                bild = shape.image
                blob = bild.blob
                ext = (bild.ext or "png").lower()

                dateiname = f"pptx_slide{i}_img{bild_zaehler}.{ext}"
                dateipfad = os.path.join(bilder_ordner, dateiname)

                with Image.open(io.BytesIO(blob)) as im:
                    im.save(dateipfad)

                folien_daten["elemente"].append(
                    {
                        "typ": "bild",
                        "datei": dateiname,
                        "z_index": z_index,
                        "position": {
                            "links": shape.left,
                            "oben": shape.top,
                            "breite": shape.width,
                            "höhe": shape.height,
                        },
                    }
                )
                continue

            # Tabellen werden als Liste von Zeilen und Zellen gespeichert
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tabellen_daten = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]

                folien_daten["elemente"].append(
                    {
                        "typ": "tabelle",
                        "inhalt": tabellen_daten,
                        "z_index": z_index,
                        "position": {
                            "links": shape.left,
                            "oben": shape.top,
                            "breite": shape.width,
                            "höhe": shape.height,
                        },
                    }
                )
                continue

            # Diagramme werden mit Typ, Kategorien und Serienwerten gespeichert
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart

                chart_data = {
                    "chart_type": str(chart.chart_type),
                    "serien": [],
                    "kategorien": None,
                }

                try:
                    plot = chart.plots[0]
                    if hasattr(plot, "categories") and plot.categories is not None:
                        chart_data["kategorien"] = [c.label for c in plot.categories]
                except Exception:
                    chart_data["kategorien"] = None

                for series in chart.series:
                    try:
                        values = list(series.values)
                    except Exception:
                        values = None
                    chart_data["serien"].append({"name": series.name, "werte": values})

                folien_daten["elemente"].append(
                    {
                        "typ": "diagramm",
                        "daten": chart_data,
                        "z_index": z_index,
                        "position": {
                            "links": shape.left,
                            "oben": shape.top,
                            "breite": shape.width,
                            "höhe": shape.height,
                        },
                    }
                )
                continue

            # Gruppenobjekte enthalten mehrere Unterelemente
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                gruppen_elemente = []
                group_z = 0

                for subshape in shape.shapes:
                    group_z += 1
                    sub_entry = {
                        "typ": str(subshape.shape_type),
                        "z_index_in_group": group_z,
                        "position": {
                            "links": subshape.left,
                            "oben": subshape.top,
                            "breite": subshape.width,
                            "höhe": subshape.height,
                        },
                    }

                    # Falls das Unterelement Text enthält, werden auch die Textdetails gespeichert
                    if getattr(subshape, "has_text_frame", False):
                        sub_tf = subshape.text_frame
                        sub_text = sub_tf.text.strip()
                        if sub_text:
                            sub_entry["inhalt"] = sub_text

                            sub_paragraphs = []
                            for p in sub_tf.paragraphs:
                                if not (p.text or "").strip():
                                    continue
                                runs_out = []
                                for r in p.runs:
                                    if not (r.text or ""):
                                        continue
                                    runs_out.append(
                                        {
                                            "text": r.text,
                                            "schriftgröße": (
                                                r.font.size.pt
                                                if (r.font and r.font.size)
                                                else None
                                            ),
                                            "schriftart": (
                                                r.font.name
                                                if (r.font and r.font.name)
                                                else None
                                            ),
                                            "farbe": safe_pptx_rgb(r),
                                            "bold": (
                                                bool(r.font.bold)
                                                if (r.font and r.font.bold is not None)
                                                else None
                                            ),
                                            "italic": (
                                                bool(r.font.italic)
                                                if (
                                                    r.font and r.font.italic is not None
                                                )
                                                else None
                                            ),
                                            "underline": (
                                                bool(r.font.underline)
                                                if (
                                                    r.font
                                                    and r.font.underline is not None
                                                )
                                                else None
                                            ),
                                        }
                                    )
                                sub_paragraphs.append(
                                    {
                                        "level": p.level if hasattr(p, "level") else 0,
                                        "is_bullet": pptx_paragraph_is_bullet(p),
                                        "textausrichtung": pptx_align_to_css(
                                            p.alignment
                                        ),
                                        "runs": runs_out,
                                        "plain_text": p.text or "",
                                    }
                                )
                            if sub_paragraphs:
                                sub_entry["paragraphs"] = sub_paragraphs

                    gruppen_elemente.append(sub_entry)

                folien_daten["elemente"].append(
                    {
                        "typ": "gruppe",
                        "anzahl_elemente": len(gruppen_elemente),
                        "elemente": gruppen_elemente,
                        "z_index": z_index,
                        "position": {
                            "links": shape.left,
                            "oben": shape.top,
                            "breite": shape.width,
                            "höhe": shape.height,
                        },
                    }
                )
                continue

            # Allgemeine Textfelder werden hier verarbeitet
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                full_text = tf.text.strip()
                if not full_text:
                    continue

                paragraphs_out = []
                for p in tf.paragraphs:
                    p_txt = (p.text or "").strip()
                    if not p_txt:
                        continue

                    runs_out = []
                    for r in p.runs:
                        t = r.text or ""
                        if not t:
                            continue
                        runs_out.append(
                            {
                                "text": t,
                                "schriftgröße": (
                                    r.font.size.pt if (r.font and r.font.size) else None
                                ),
                                "schriftart": (
                                    r.font.name if (r.font and r.font.name) else None
                                ),
                                "farbe": safe_pptx_rgb(r),
                                "bold": (
                                    bool(r.font.bold)
                                    if (r.font and r.font.bold is not None)
                                    else None
                                ),
                                "italic": (
                                    bool(r.font.italic)
                                    if (r.font and r.font.italic is not None)
                                    else None
                                ),
                                "underline": (
                                    bool(r.font.underline)
                                    if (r.font and r.font.underline is not None)
                                    else None
                                ),
                            }
                        )

                    paragraphs_out.append(
                        {
                            "level": p.level if hasattr(p, "level") else 0,
                            "is_bullet": pptx_paragraph_is_bullet(p),
                            "textausrichtung": pptx_align_to_css(p.alignment),
                            "zeilenabstand": p.line_spacing,
                            "abstand_vor": (
                                p.space_before.pt if p.space_before else None
                            ),
                            "abstand_nach": (
                                p.space_after.pt if p.space_after else None
                            ),
                            "runs": runs_out,
                            "plain_text": p.text or "",
                        }
                    )

                # Falls keine Absatzstruktur erkannt wurde, wird ein einfacher Fallback angelegt
                if not paragraphs_out:
                    paragraphs_out = [{"level": 0, "runs": [{"text": full_text}]}]

                folien_daten["elemente"].append(
                    {
                        "typ": "text",
                        "inhalt": full_text,
                        "paragraphs": paragraphs_out,
                        "schriftgröße": None,
                        "schriftart": None,
                        "farbe": None,
                        "textausrichtung": None,
                        "zeilenabstand": None,
                        "abstand_vor": None,
                        "abstand_nach": None,
                        "listebene": 0,
                        "z_index": z_index,
                        "position": {
                            "links": shape.left,
                            "oben": shape.top,
                            "breite": shape.width,
                            "höhe": shape.height,
                        },
                    }
                )
        # JSON-Datei für die aktuelle Folie speichern
        json_name = f"pptx_slide_{i:03d}.json"
        json_path = os.path.join("data/processed", json_name)

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(folien_daten, f, ensure_ascii=False, indent=2)

        audit_logger.info(f"{json_name} | Hash: {hash_file(json_path)}")

    runtime = time.time() - start_time
    meta = {
        "python_version": sys.version,
        "pptx_version": pptx.__version__,
        "pymupdf_version": fitz.__doc__,
        "system": platform.platform(),
        "runtime_seconds": runtime,
    }

    # Metadaten zur Laufzeit und Umgebung speichern
    with open("data/processed/pipeline_meta.json", "w", encoding="utf-8") as f:
        json.dump(meta, f, indent=2, ensure_ascii=False)

    logging.info("PPTX Parsing abgeschlossen")
    audit_logger.info("PPTX Parsing abgeschlossen")
    
    # Ausgabe-Dateien und Assets nochmals mit Hash protokollieren
    audit_hash_artifacts("data/processed", "PHASE1_OUTPUT_PROCESSED")
    audit_hash_artifacts("export/assets", "PHASE1_OUTPUT_ASSETS")


# =================================================================================================================
# PDF PARSING: Verarbeitet eine PDF-Datei und exportiert Text, Bilder und erkannte Diagrammregionen als JSON/Assets
# =================================================================================================================
def parse_pdf(pdf_path: str) -> None:
    logging.info(f"PDF Parsing gestartet: {pdf_path}")
    start_time = time.time()

    input_hash = hash_file(pdf_path)
    audit_logger.info(f"PDF Input: {pdf_path} | Hash: {input_hash}")

    doc = fitz.open(pdf_path)

    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        rect = page.rect

        folien_daten = {
            "foliennummer": page_num + 1,
            "überschrift": "kein Titel gefunden",
            "canvas": {
                "width": float(rect.width),
                "height": float(rect.height),
                "unit": "pt",
            },
            "elemente": [],
        }

        z_index = 0

        # Bilder der PDF-Seite extrahieren
        for img_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base_img = doc.extract_image(xref)
            img_bytes = base_img["image"]
            img_ext = base_img["ext"]

            filename = f"pdf_page{page_num+1}_img{img_index}.{img_ext}"
            filepath = os.path.join("export/assets", filename)

            with open(filepath, "wb") as f:
                f.write(img_bytes)

            rects = page.get_image_rects(xref)
            if rects:
                position = [
                    {
                        "links": float(r.x0),
                        "oben": float(r.y0),
                        "breite": float(r.x1 - r.x0),
                        "höhe": float(r.y1 - r.y0),
                    }
                    for r in rects
                ]
            else:
                position = [{"links": None, "oben": None, "breite": None, "höhe": None}]

            z_index += 1
            folien_daten["elemente"].append(
                {
                    "typ": "bild",
                    "datei": filename,
                    "z_index": z_index,
                    "position": position,
                }
            )

        # Textblöcke auslesen
        text_blocks = extract_pdf_text_blocks(page)

        # Titel automatisch aus dem Text bestimmen
        title = detect_pdf_title_from_textblocks(
            text_blocks, page_height=float(rect.height)
        )
        if title:
            folien_daten["überschrift"] = title

        # Textblöcke mit z-Index übernehmen
        for el in text_blocks:
            z_index += 1
            el["z_index"] = z_index
            folien_daten["elemente"].append(el)

        # Gezeichnete Elemente analysieren, um Diagrammbereiche zu finden
        drawings = page.get_drawings()
        chart_regions = detect_chart_regions(page, folien_daten["elemente"], drawings)

        chart_img_idx = 0
        for reg in chart_regions:
            z_index += 1
            chart_img_idx += 1

            x0, y0, x1, y1 = reg["bbox"]
            if None in (x0, y0, x1, y1):
                continue
            if x1 <= x0 or y1 <= y0:
                continue

            chart_filename = f"pdf_page{page_num+1}_chart{chart_img_idx}.png"
            chart_path = os.path.join("export/assets", chart_filename)

            # Diagrammbereich wird ausgeschnitten und als Bild gespeichert
            try:
                clip = fitz.Rect(float(x0), float(y0), float(x1), float(y1))
                pix = page.get_pixmap(clip=clip, dpi=200, alpha=False)
                pix.save(chart_path)
            except Exception as e:
                logging.warning(f"Chart-Crop fehlgeschlagen Seite {page_num+1}: {e}")
                chart_filename = None

            folien_daten["elemente"].append(
                {
                    "typ": "diagramm_region",
                    "bbox": {
                        "links": float(x0),
                        "oben": float(y0),
                        "breite": float(x1 - x0),
                        "höhe": float(y1 - y0),
                    },
                    "labels": reg.get("labels", []),
                    "chart_image": chart_filename,
                    "z_index": z_index,
                }
            )
        # JSON-Datei für die PDF-Seite speichern
        json_name = f"pdf_slide_{page_num+1:03d}.json"
        json_path = os.path.join("data/processed", json_name)

        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(folien_daten, f, ensure_ascii=False, indent=2)

        audit_logger.info(f"{json_name} | Hash: {hash_file(json_path)}")

    doc.close()

    runtime = time.time() - start_time
    logging.info(f"PDF Processing Time: {runtime:.2f}s")
    audit_logger.info("PDF Parsing abgeschlossen")

    # Alle erzeugten Dateien und Assets mit Hash protokollieren
    audit_hash_artifacts("data/processed", "PHASE1_OUTPUT_PROCESSED")
    audit_hash_artifacts("export/assets", "PHASE1_OUTPUT_ASSETS")


# ==========================================================
# AUSFÜHRUNG
# ==========================================================
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--pptx",
        default=r"C:\Users\oezdemir\Desktop\Projekt\data\raw\slides_input\Test1_neu.pptx",
    )
    parser.add_argument(
        "--pdf",
        default=r"C:\Users\oezdemir\Desktop\Projekt\data\raw\slides_input\TEST2_neu1.pdf",
    )
    args = parser.parse_args()

    # Beide Parser nacheinander ausführen
    parse_pptx(args.pptx)
    parse_pdf(args.pdf)

    # Erzeugt eine Mapping-Datei mit Eingaben und Ausgaben
    mapping = {
        "pptx_input": args.pptx,
        "pdf_input": args.pdf,
        "outputs": sorted(os.listdir("data/processed")),
    }

    map_path = "data/processed/map.json"
    with open(map_path, "w", encoding="utf-8") as f:
        json.dump(mapping, f, indent=2, ensure_ascii=False)

    audit_logger.info(f"map.json | Hash: {hash_file(map_path)}")
    print("Pipeline abgeschlossen - Logs & Hash erstellt")


if __name__ == "__main__":
    main()